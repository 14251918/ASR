import os
import json
import base64
import hashlib
import asyncio
import httpx 
from dotenv import load_dotenv

from tenacity import retry, stop_after_attempt, wait_fixed, retry_if_exception_type

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import PromptTemplate
from langchain_mistralai.chat_models import ChatMistralAI
from langchain_core.messages import HumanMessage

load_dotenv()
os.environ["MISTRAL_API_KEY"] = os.getenv("MISTRAL_API_KEY")

model_context = ChatMistralAI(
    model="pixtral-large-latest",
    temperature=0.0
)

def compute_slide_hash(slide):
    hasher = hashlib.sha256()
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            hasher.update(shape.text.strip().encode("utf-8"))
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        hasher.update(cell.text.strip().encode("utf-8"))
        elif hasattr(shape, "_element") and shape._element.tag.endswith("graphicFrame"):
            for node in shape._element.findall(".//a:t", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}):
                text = node.text
                if text and text.strip():
                    hasher.update(text.strip().encode("utf-8"))

    def hash_shape(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            hasher.update(shape.image.blob)
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image"):
            hasher.update(shape.image.blob)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                hash_shape(sub_shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.text.strip():
            hasher.update(shape.text.strip().encode("utf-8"))
    for shape in slide.shapes:
        hash_shape(shape)
    return hasher.hexdigest()

@retry(
    stop=stop_after_attempt(5),
    wait=wait_fixed(1),
    retry=retry_if_exception_type(httpx.HTTPStatusError),
    retry_error_callback=lambda retry_state: f"Не удалось обработать после {retry_state.attempt_number} попыток"
)

async def process_image_with_retry(model, message):
    response = await model_context.ainvoke([message])
    return response.content

@retry(
    stop=stop_after_attempt(5),
    wait=wait_fixed(1),
    retry=retry_if_exception_type(httpx.HTTPStatusError),
    retry_error_callback=lambda retry_state: {
        "text_from_image": "Ошибка обработки",
        "description": f"Не удалось обработать после {retry_state.attempt_number} попыток"
    }
)
def process_json_with_retry(json_data):
    prompt_template = PromptTemplate(
        input_variables=["json_data"],
        template="""
            Ты — ассистент, который анализирует JSON с данными о слайдах презентации. Для каждого слайда предоставь краткий контекст (назначение слайда, основная идея или содержание) на основе текста и описаний изображений. Если данные из изображений отсутствуют, предположи контекст на основе заголовков и структуры презентации. JSON:

                {json_data}

                Формат ответа:
                [
                    {{
                        "slide_number": <номер слайда>,
                        "context": "<краткий контекст слайда>"
                    }},
                    ...
                ]
                """
    )
    json_str = json.dumps(json_data, ensure_ascii=False, indent=4)
    chain = prompt_template | model_context | StrOutputParser()
    description = chain.invoke({"json_data": json_str})
    return description

def clean_text_from_image(text):
    if not text:
        return "Отсутствует"
    lines = list(dict.fromkeys(line.strip() for line in text.splitlines() if line.strip()))
    return "\n".join(lines) if lines else "Отсутствует"

def extract_smartart_text(shape):
    try:
        xml_bytes = etree.tostring(shape._element, encoding="utf-8", xml_declaration=False)
    except Exception as e:
        print(f"Не удалось извлечь объекты из SmartArt {shape.name}: {e}")
        return []

    parser = etree.XMLParser(recover=True)
    try:
        root = etree.fromstring(xml_bytes, parser)
    except Exception as e:
        print(f"Не удалось извлечь текст из SmartArt {shape.name}: {e}")
        return []

    ns = {"dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram"}
    rel_ids = root.find(".//dgm:relIds", ns)
    if rel_ids is None:
        print(f"В SmartArt {shape.name} нет узлов")
        return []
    texts = []
    for attr, rel_id in rel_ids.attrib.items():
        try:
            rel = shape.part.rels[rel_id]
            part = rel.target_part
            blob = part.blob
            droot = etree.fromstring(blob, parser)

            for t in droot.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                raw = t.text or ""
                text = raw.strip()
                if text:
                    texts.append(text)
        except KeyError:
            print(f"Ошибка при обработке {shape.name}")
            continue
    return texts

def extract_images_from_shape(shape, images, group_text, slide_num):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        if len(shape.image.blob) < 1000:
            print(f"Изображение (picture) {shape.name} слишком маленькое")
            return
        print(f"Найдено изображение (picture): {shape.name}")
        images.append((shape.image.blob, shape.image.ext.lower(), shape.name, slide_num))
    elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        if hasattr(shape, "image"):
            if len(shape.image.blob) < 1000:
                print(f"Изображение (placeholder) {shape.name} слишком маленькое")
                return
            print(f"Найдено изображение (placeholder): {shape.name}")
            images.append((shape.image.blob, shape.image.ext.lower(), shape.name, slide_num))
        elif hasattr(shape, "text") and shape.text.strip():
            group_text.append(shape.text.strip())
        else:
            print(f"{shape.name} (placeholder) без изображения или текста")
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"Найдена группа фигур: {shape.name}")
        for sub_shape in shape.shapes:
            extract_images_from_shape(sub_shape, images, group_text, slide_num)
    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.text.strip():
        print(f"Найден текст: {shape.text.strip()}")
        group_text.append(shape.text.strip())
    else:
        smartart_text = extract_smartart_text(shape)
        if smartart_text:
            group_text.extend(smartart_text)
        else:
            print(f"Пропущено: {shape.name} типа {shape.shape_type}")

async def process_presentation(file_path):
    if not os.path.exists(file_path):
        print("Такого файла не существует")
        return
    try:
        context_file_name = os.path.splitext(os.path.basename(file_path))[0]
        output_context_json = f"{context_file_name}.json"

        pres = Presentation(file_path)
        slides = list(pres.slides)
        total_slides = len(slides)
        print(f"Всего слайдов: {total_slides}")

        if os.path.exists(output_context_json):
            with open(output_context_json, "r", encoding="utf-8") as f:
                output_data = json.load(f)
            print(f"Найден существующий JSON с {len(output_data)} обработанными слайдами")
            if total_slides == len(output_data):
                return output_data
        else:
            output_data = []
            print("Создание JSON файла")
        for slide_num, slide in enumerate(slides, 1):
            cur_hash = compute_slide_hash(slide)
            slide_data = next((data for data in output_data if data["slide_number"] == slide_num), None)
            if slide_data and slide_data.get("hash") == cur_hash:
                print(f"Слайд {slide_num}, хэш совпадает")
                continue
            else:
                print(f"Слайд {slide_num} обрабатывается")

            slide_data = {
                "slide_number": slide_num,
                "hash": cur_hash,
                "text": [],
                "images": [],
                "group_text": []
            }

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_data["text"].append(shape.text.strip())
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_text = []
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                table_text.append(cell.text.strip())
                    if table_text:
                        slide_data["text"].extend(table_text)
            images = []
            group_text = []
            print(f"Слайд {slide_num}: поиск изображений")
            for shape in slide.shapes:
                extract_images_from_shape(shape, images, group_text, slide_num)

            print(f"Слайд {slide_num}: найдено {len(images)} изображений")
            if group_text:
                slide_data["group_text"] = group_text
                print(f"Слайд {slide_num}: текст: {group_text}")

            tasks = []
            for i, (image_bytes, image_ext, image_name, image_num) in enumerate(images):
                image_base64 = base64.b64encode(image_bytes).decode("utf-8")
                mime_type = "data:image/png;base64"
                message = HumanMessage(
                    content=[
                        {
                            "type": "text",
                            "text": """
                            Ты — ассистент, который анализирует изображение. Извлеки текст (если он есть) и опиши содержимое изображения.

                            Формат ответа:
                            - Текст из изображения: <текст или "Отсутствует">
                            - Описание изображения: <описание или "Отсутствует">
                            """
                        },
                        {
                            "type": "image_url",
                            "image_url": f"{mime_type},{image_base64}"
                        }
                    ]
                )
                tasks.append(process_image_with_retry(model_context, message))
            responses = await asyncio.gather(*tasks, return_exceptions=True)
            for i, (image_bytes, image_ext, image_name, image_num) in enumerate(images):
                response = responses[i]
                if isinstance(response, str):
                    print(f"Слайд {slide_num}: изображение {image_name} обработано")
                    image_data = {
                        "image_name": image_name,
                        "text_from_image": "Отсутствует",
                        "description": "Отсутствует",
                    }
                    parsed_response = parse_model_response(response)
                    image_data["text_from_image"] = parsed_response["text_from_image"]
                    image_data["description"] = parsed_response["description"]
                    slide_data["images"].append(image_data)
                else:
                    print(f"Слайд {slide_num}: ошибка обработки изображения {image_name}: {response}")
                    image_data = {
                        "image_name": image_name,
                        "text_from_image": "Ошибка обработки",
                        "description": str(response)
                    }
                    slide_data["images"].append(image_data)
            output_data = [data for data in output_data if data["slide_number"] != slide_num]
            output_data.append(slide_data)

            output_data.sort(key=lambda x: x["slide_number"])
            with open(output_context_json, "w", encoding="utf-8") as f:
                json.dump(output_data, f, ensure_ascii=False, indent=4)
            print(f"Слайд {slide_num}: сохранен в {output_context_json}")
        try:
            response = process_json_with_retry(output_data)
            print(f"Ответ модели на JSON: {response}")
            response = response.strip()
            if response.startswith("```json"):
                response = response[7:-3].strip()
            context_data = json.loads(response)
            for slide_data in output_data:
                context = next((item["context"] for item in context_data if item["slide_number"] == slide_data["slide_number"]), "Без контекста")
                slide_data["slide_context"] = context
        except json.JSONDecodeError as e:
            print(f"Ошибка обработки JSON от модели: {e}")
        except Exception as e:
            print(f"Ошибка при обработке JSON: {e}")


        output_data.sort(key=lambda x: x["slide_number"])
        with open(output_context_json, "w", encoding="utf-8") as f:
            json.dump(output_data, f, ensure_ascii=False, indent=4)
        print(f"JSON сохранен в {output_context_json} ({total_slides} слайдов)")
        return output_data
    except Exception as e:
        print(f"Ошибка при обработке презентации: {e}")
        return None

def parse_model_response(response):
    text_from_image = []
    description = []
    current_section = None

    if not isinstance(response, str) or not response.strip():
        return {
            "text_from_image": "Отсутствует",
            "description": "Пустой ответ модели"
        }
    lines = response.split("\n")
    for line in lines:
        line = line.strip()
        if line.startswith("- Текст из изображения:"):
            current_section = "text"
            line = line.replace("- Текст из изображения:", "").strip()
            text_from_image.append(line)
        elif line.startswith("- Описание изображения:"):
            current_section = "description"
            line = line.replace("- Описание изображения:", "").strip()
            description.append(line)
        elif line and current_section == "text":
            text_from_image.append(line)
        elif line and current_section == "description":
            description.append(line)
    text = "\n".join(text_from_image).strip() if text_from_image else "Отсутствует"
    desc = "\n".join(description).strip() if description else "Отсутствует"
    return {
            "text_from_image": clean_text_from_image(text),
            "description": desc
        }
