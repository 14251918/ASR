import datetime
import json
import os
import re

from langchain_core.output_parsers import StrOutputParser
from langchain_groq import ChatGroq
from langchain.prompts import PromptTemplate
import fitz
from pptx import Presentation
import queue
import threading
import time
import numpy as np
import sounddevice as sd
import win32com.client
import io
import wave
from groq import Groq

os.environ["GROQ_API_KEY"] = "API_KEY"
file_path = r"presentation.pptx"

class ASR:
    def __init__(self,
                 file_path,
                 sample_rate=16000,
                 LLM_model_name="gemma2-9b-it"
                 ):
        self.ppAPP = None
        self.slide_contexts = None
        self.count_slides = None
        self.slideshow = None
        self.presentation_settings = None
        self.presentation = None
        self.sample_rate = sample_rate
        self.blocksize = self.sample_rate * 2
        self.min_louds = self.sample_rate * 0.1
        self.min_speach = int(self.sample_rate * 0.5)
        self.max_speech = self.sample_rate * 4
        self.mean_back = 0
        self.channels = 1
        self.dtype = "float32"
        self.buffer = np.array([])
        self.file_path = file_path
        if not os.path.isabs(file_path):
            self.file_path = os.path.abspath(self.file_path)
        self.file_path = self.file_path.replace("\\", "//")
        self.extension = os.path.splitext(self.file_path)[-1].lower()
        self.was_audio = 0
        self.was_audio_recognized = 0
        self.was_LLM_recognized = 0
        self.commands = {
            "next_slide",
            "prev_slide",
            "close_slideshow",
            "move_to_slide",
            "no_move"
        }
        self.command_queue = queue.Queue()
        self.client = Groq()
        self.model_name = LLM_model_name
        self.model = ChatGroq(model_name=self.model_name, temperature=0)
        self.output_parser = StrOutputParser()

        self.stream = sd.InputStream(
            samplerate=self.sample_rate,
            blocksize=self.blocksize,
            channels=self.channels,
            callback=self.audio_callback,
            dtype=self.dtype
        )



        self.template_to_command = '''
Ты - ассистент-помощник для переключения слайдов.
Дана распознанная речь, возможно, с ошибками: {recognized_speech}. 

Текущий слайд: {slide_now}.
Контекст слайдов (в формате JSON): """{slide_context}""" Контекст закончен.

Твоя задача - вывести функцию для управления слайдами на основе распознанной речи.
Доступные функции:
- next_slide: Перейти на следующий слайд.
- prev_slide: Перейти на предыдущий слайд.
- close_slideshow: Закрыть презентацию.
- move_to_slide(num_slide): Перейти на слайд с номером num_slide.
- no_move: Ничего не делать, если нет команды управления.

Важно:
- Используй next_slide, если сказано "следующий", "дальше", "перейти на следующий слайд" и т.п.
- Используй prev_slide, если сказано "предыдущий", "назад", "вернуться" и т.п.
- Используй move_to_slide(num_slide), если указан конкретный номер слайда, например, "перейти на слайд 5", "десятый слайд".
Для последнего или первого слайда также используются числа
- Если речь не содержит явной команды управления слайдами, используй no_move.
Без додумывания. Пользователь должен  обратиться к тебе для переключения, или в речи есть слова, подходящие к контексту слайда.
Твой ответ должен быть в формате:
functions_name(args) — если есть аргументы,
или просто functions_name — если аргументов нет.

**Не добавляй никаких дополнительных символов, тегов (например, ```python) или объяснений в первой строке ответа.**
Объяснение пиши после пустой строки.

Примеры:
- Распознанная речь: "следующий" → next_slide
- Распознанная речь: "следующий слайд" → next_slide
- Распознанная речь: "перейти на слайд 3" → move_to_slide(3)
- Распознанная речь: "десятый слайд" → move_to_slide(10)
- Распознанная речь: "какой сейчас слайд" → no_move

Учти, что это примеры.
Это не распознанная речь. Распознанная речь - в начале.
Также учти, что ошибки в словах могут быть значительны. 
Например, слово слойт - может быть словом слайд и т.п.

'''

        self.prompt_command = PromptTemplate(
            input_variables=["slide_now",
                             "slide_context",
                             "recognized_speech",
                             "functions_name"],
            template=self.template_to_command
        )
        self.chain_command = (self.prompt_command |
                              self.model |
                              self.output_parser)

    def run_powerPoint(self, file_path=""):
        if file_path:
            self.file_path = file_path
        if not os.path.exists(self.file_path):
            print(f"Такого файла {self.file_path} не существует")
            return
        try:
            self.ppAPP = win32com.client.GetActiveObject("PowerPoint.Application")
        except Exception as e:
            self.ppAPP = win32com.client.Dispatch("PowerPoint.Application")
            self.ppAPP.Visible = True
        try:
            self.presentation = self.ppAPP.Presentations.Open(self.file_path)
            self.presentation_settings = self.presentation.SlideShowSettings
            self.presentation_settings.Run()
            for _ in range(5):
                if self.ppAPP.SlideShowWindows.Count:
                    self.slideshow = self.ppAPP.SlideShowWindows(1)
                    break
                time.sleep(1)
            self.count_slides = self.presentation.Slides.Count

        except Exception as e:
            print(f"Ошибка: {e}. Пожалуйста, укажите полный путь к презентации.")
            if self.ppAPP:
                self.ppAPP.Quit()

    def next_slide(self):
        if self.count_slides > self.slideshow.View.Slide.SlideIndex:
            self.slideshow.View.Next()
        else:
            print("Это последний слайд.")

    def prev_slide(self):
        if self.slideshow.View.Slide.SlideIndex > 1:
            self.slideshow.View.Previous()
        else:
            print("Это первый слайд.")

    def move_to_slide(self, slide_number):
        try:
            slide_number = int(slide_number)
            if 1 <= slide_number <= self.count_slides:
                self.slideshow.View.GotoSlide(slide_number)
            else:
                print(f"Недопустимый номер слайда {slide_number}.")
        except ValueError:
            print("Недопустимый номер слайда. Ожидалось число.")

    def close_slideshow(self):
        if self.slideshow:
            self.slideshow.View.Exit()

    def no_move(self):
        print("Ничего не делать. Выполнено!")

    def __exit__(self):
        if self.ppAPP:
            self.ppAPP.DisplayAlerts = False
            self.ppAPP.Quit()
        del self.ppAPP

    def asr(self, audio):
        def recognize(audio_to_rec):
            if np.max(np.abs(audio_to_rec)) != 0:
                audio_to_rec = audio_to_rec / np.max(np.abs(audio_to_rec))
            audio_to_rec = np.int16(audio_to_rec * 32767)
            with io.BytesIO() as wav_buffer:
                with wave.open(wav_buffer, 'wb') as wf:
                    wf.setnchannels(1)
                    wf.setsampwidth(2)
                    wf.setframerate(self.sample_rate)
                    wf.writeframes(audio_to_rec.tobytes())
                wav_buffer.seek(0)

                recognized_text = self.client.audio.transcriptions.create(
                    file=("recording.wav", wav_buffer),
                    # model="whisper-large-v3",
                    model="whisper-large-v3-turbo",
                    language="ru",
                    response_format="text"
                )
            self.was_audio_recognized += 1
            print(f"Аудио №{self.was_audio_recognized} распознано в {datetime.datetime.now()}.")
            print(f"Распознанная речь : {recognized_text}")
            self.command_queue.put(recognized_text.strip().lower())

        threading.Thread(target=recognize, args=(audio,)).start()

    def audio_callback(self, indata, frames, time, status, alpha=0.2):
        indata = indata.squeeze()
        mean = np.mean(np.abs(indata))
        self.mean_back = alpha * mean + (1 - alpha) * self.mean_back

        loud_frames = sum(abs(i) > self.mean_back * 2 + 0.001 for i in indata)

        def start_recognize():
            self.was_audio += 1
            print(f"Обработка аудио №{self.was_audio} начата в {datetime.datetime.now()}")
            self.asr(self.buffer.copy())
            self.buffer = self.buffer[-self.min_speach:]

        if loud_frames > self.min_louds and len(self.buffer) < self.max_speech:
            print(loud_frames, self.mean_back)
            self.buffer = np.append(self.buffer, indata)
            if len(self.buffer) >= self.max_speech:
                start_recognize()
        elif len(self.buffer) > self.min_speach:
            start_recognize()
        else:
            self.buffer = np.array([])


    def process_command(self, recognized_speech):
        try:
            self.was_LLM_recognized += 1
            num = self.was_LLM_recognized
            print(f"Обработка команды №{num} начата в {datetime.datetime.now()}.")
            slide_now = self.slideshow.View.Slide.SlideIndex
            functions_name = ', '.join(self.commands)
            now_time = datetime.datetime.now()
            output = self.chain_command.invoke(
                {"slide_now": slide_now,
                 "slide_context": json.dumps(self.slide_contexts, ensure_ascii=False),
                 "recognized_speech": recognized_speech,
                 "functions_name": functions_name}
            )
            print(f"Команда №{num} обработана в {datetime.datetime.now()}.")
            print(f"Распознанный текст: {recognized_speech}, Ответ LLM: {output}")
            match = re.search(r'(\w+)(?:\((\d+)\))?', output)
            if match:
                fun_name = match.group(1)
                args_str = match.group(2) if match.group(2) else None
                print(f"Команда: {fun_name}, аргументы: {args_str}")
            else:
                print("Команда не найдена в ответе LLM")
                return

            if fun_name in self.commands:
                action = fun_name
                print(f"Найдена команда: {action}")
                try:
                    if args_str:
                        if args_str.isdigit():
                            getattr(self, action)(int(args_str))
                        else:
                            print(f"Недопустимый аргумент: {args_str}")
                    else:
                        getattr(self, action)()
                except Exception as e:
                    print(f"Ошибка при выполнении команды: {e}")
            else:
                print(f"Неизвестная команда {fun_name}")
        except Exception as e:
            print(f"Ошибка при обработке распознанной речи: {e}")


    def start_audio(self):
        with self.stream:
            print("Запись начата")
            try:
                while True:
                    try:
                        recognized_speech = self.command_queue.get_nowait()
                        print("Команда обрабатывается.")
                        self.process_command(recognized_speech)
                    except queue.Empty:
                        time.sleep(0.1)
            except KeyboardInterrupt:
                print("Запись остановлена")
                self.__exit__()


    def read_file_pdf(self, pdf_path=''):
        if pdf_path:
            self.file_path = pdf_path
            self.extension = os.path.splitext(self.file_path)[-1].lower()
        self.slide_contexts = {}
        try:
            pdf_doc = fitz.open(self.file_path)
            for page_num in range(pdf_doc.page_count):
                page = pdf_doc.load_page(page_num)
                text = page.get_textpage()
                self.slide_contexts[str(page_num + 1)] = text.extractText()
            pdf_doc.close()
        except Exception as e:
            print(f"Ошибка извлечения из PDF: {e}")

    def read_file_pptx(self, pptx_path=''):
        if pptx_path:
            self.file_path = pptx_path
            self.extension = os.path.splitext(self.file_path)[-1].lower()
        self.slide_contexts = {}
        try:
            present = Presentation(self.file_path)
            for i, slide in enumerate(present.slides):
                text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                self.slide_contexts[str(i + 1)] = text.strip()
        except Exception as e:
            print(f"Ошибка извлечения из PPTX: {e}")


    def read_file(self, file_path=''):
        if file_path:
            self.file_path = file_path
            self.extension = os.path.splitext(self.file_path)[-1].lower()
        if self.extension == ".pdf":
            self.read_file_pdf(self.file_path)
        elif self.extension == ".pptx":
            self.read_file_pptx(self.file_path)
        else:
            print("Данный тип файла не поддерживается.")

model = ASR(file_path=file_path)
model.run_powerPoint()
model.read_file()
model.start_audio()
